from pptx import Presentation
from pptx.util import Emu
import sys

path = "/home/norepinephrine/Documents/Heart-Conduction/MonthlyReport/Chang - Boundary Handling in Discrete Cardiac Diffusion.pptx"
prs = Presentation(path)
print(f"slide size: {Emu(prs.slide_width).inches:.2f} x {Emu(prs.slide_height).inches:.2f} in")
print(f"n slides: {len(prs.slides)}")
print("=" * 70)
for i, slide in enumerate(prs.slides, 1):
    layout = slide.slide_layout.name
    n_pics = sum(1 for sh in slide.shapes if sh.shape_type == 13)  # PICTURE
    n_movies = sum(1 for sh in slide.shapes if sh.shape_type == 16)  # MEDIA
    texts = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip().replace("\n", " / ")
            if t:
                texts.append(t)
    title = texts[0][:80] if texts else "(no text)"
    print(f"\n[Slide {i}]  layout='{layout}'  pics={n_pics} media={n_movies}")
    print(f"  TITLE: {title}")
    for t in texts[1:6]:
        print(f"   - {t[:110]}")
