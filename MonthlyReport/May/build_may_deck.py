"""Build the May progress-report deck from the boundary talk.
Adds two slides (BC-mechanism summary, textbook collage), each with a title and
a grey takeaway box, then positions them. Original talk file is left untouched.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image
import copy

SRC = "/home/norepinephrine/Documents/Heart-Conduction/MonthlyReport/Chang - Boundary Handling in Discrete Cardiac Diffusion.pptx"
OUT = "/home/norepinephrine/Documents/Heart-Conduction/MonthlyReport/May/Chang - Progress Report - May 2026.pptx"
FIG = "/home/norepinephrine/Documents/Heart-Conduction/MonthlyReport/May/figures"

NAVY = RGBColor(0x1F, 0x3C, 0x88)
GREY_FILL = RGBColor(0xEC, 0xEC, 0xEC)
DARK = RGBColor(0x22, 0x22, 0x22)

prs = Presentation(SRC)
SW, SH = prs.slide_width, prs.slide_height  # EMU

# reuse the exact layout the content slides use (slide 8 = "Streaming vs Receiving")
LAYOUT = prs.slides[7].slide_layout
print("using layout:", LAYOUT.name, "placeholders:", len(LAYOUT.placeholders))

WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def set_title(slide, text):
    """Replicate the deck's heading textbox (centered 28pt bold white on the master's
    charcoal bar). The deck uses a manual textbox here, not a title placeholder."""
    box = slide.shapes.add_textbox(Inches(0.24), Inches(0.16), Inches(9.51), Inches(0.67))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = WHITE
    return box

def add_image_fit(slide, path, top_in, max_w_in, max_h_in):
    iw, ih = Image.open(path).size
    ar = iw / ih
    w = max_w_in; h = w / ar
    if h > max_h_in:
        h = max_h_in; w = h * ar
    left = (Emu(SW).inches - w) / 2.0
    slide.shapes.add_picture(path, Inches(left), Inches(top_in), Inches(w), Inches(h))
    return top_in + h

def add_takeaway(slide, text):
    h = 0.92
    top = Emu(SH).inches - h - 0.18
    box = slide.shapes.add_textbox(Inches(0.4), Inches(top), SW - Inches(0.8), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = GREY_FILL
    box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC); box.line.width = Pt(0.75)
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Takeaway:  "
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = DARK; r.font.name = "Calibri"
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(12); r2.font.color.rgb = DARK; r2.font.name = "Calibri"

def new_slide(title, img, takeaway):
    s = prs.slides.add_slide(LAYOUT)
    set_title(s, title)
    add_image_fit(s, img, top_in=1.10, max_w_in=9.3, max_h_in=5.15)
    add_takeaway(s, takeaway)
    return s

new_slide(
    "Boundary Rules Set the Crescent",
    f"{FIG}/bc_mechanism.png",
    "one source–sink imbalance, three wall rules — face-mirror / HBB starves the edge "
    "(forward crescent), specular is neutral, and the new horizontal-redirect over-feeds it "
    "(inverse crescent = boundary speed-up).",
)
new_slide(
    "A Visual Cardiac-Modeling Textbook",
    f"{FIG}/textbook_collage.png",
    "a self-authored 139-page visual textbook (ion channels → bidomain solvers, equations "
    "verified against Engine V5.4) is the conceptual backbone behind the project's engines.",
)

# ---- reorder: put the two new slides just before the final 'Future Goals' slide ----
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
n = len(ids)
bc_id, tb_id = ids[n - 2], ids[n - 1]   # the two we just appended
# original last slide (Future Goals) is at index n-3 now
fg_index = n - 3
sldIdLst.remove(bc_id); sldIdLst.remove(tb_id)
ids = list(sldIdLst)
# insert bc then tb right before Future Goals (which is at fg_index in the trimmed list)
sldIdLst.insert(fg_index, tb_id)
sldIdLst.insert(fg_index, bc_id)

prs.save(OUT)
print("saved", OUT)
print("final slide count:", len(prs.slides))
